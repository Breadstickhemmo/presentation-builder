from functools import wraps
from flask import request, jsonify, Blueprint, current_app, g, send_file
import jwt
import io
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt
from ..models import User, Presentation, Slide, SlideElement
from ..extensions import db

presentations_bp = Blueprint('presentations', __name__)

PIXELS_PER_INCH = 96.0

def px_to_inches(px):
    return px / PIXELS_PER_INCH

def token_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        token = None
        if 'authorization' in request.headers: token = request.headers['authorization'].split(' ')[1]
        if not token: return jsonify({'message': 'Токен аутентификации отсутствует'}), 401
        try:
            data = jwt.decode(token, current_app.config['SECRET_KEY'], algorithms=["HS256"])
            g.current_user = User.query.get(data['user_id'])
            if not g.current_user: return jsonify({'message': 'Пользователь не найден'}), 401
        except: return jsonify({'message': 'Недействительный токен'}), 401
        return f(*args, **kwargs)
    return decorated

@presentations_bp.route('/presentations/<string:presentation_id>/download', methods=['GET'])
@token_required
def download_presentation(presentation_id):
    presentation = Presentation.query.get_or_404(presentation_id)
    if presentation.user_id != g.current_user.id: return jsonify({'message': 'Доступ запрещен'}), 403
    
    prs = PptxPresentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    slides = Slide.query.filter_by(presentation_id=presentation.id).order_by(Slide.slide_number).all()

    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for element in slide_data.elements:
            if element.element_type == 'TEXT':
                left = Inches(px_to_inches(element.pos_x))
                top = Inches(px_to_inches(element.pos_y))
                width = Inches(px_to_inches(element.width))
                height = Inches(px_to_inches(element.height))
                
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = element.content or ""
                tf.word_wrap = True
                if tf.paragraphs:
                    tf.paragraphs[0].font.size = Pt(element.font_size or 24)

    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    return send_file(file_stream, as_attachment=True, download_name=f"{presentation.title}.pptx", mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

@presentations_bp.route('/presentations', methods=['POST'])
@token_required
def create_presentation():
    data = request.get_json()
    title = data.get('title', 'Новая презентация')
    new_presentation = Presentation(title=title, owner=g.current_user)
    db.session.add(new_presentation)
    db.session.flush()

    first_slide = Slide(slide_number=1, presentation_id=new_presentation.id)
    db.session.add(first_slide)
    db.session.flush()

    title_element = SlideElement(slide_id=first_slide.id, element_type='TEXT', content='Ваш заголовок', pos_x=100, pos_y=100, width=1080, height=150, font_size=44)
    subtitle_element = SlideElement(slide_id=first_slide.id, element_type='TEXT', content='Ваш подзаголовок', pos_x=100, pos_y=260, width=1080, height=100, font_size=28)
    db.session.add(title_element)
    db.session.add(subtitle_element)
    
    db.session.commit()

    first_slide_data = {
        'id': first_slide.id,
        'slide_number': first_slide.slide_number,
        'background_color': first_slide.background_color,
        'elements': [
            {'id': e.id, 'element_type': e.element_type, 'pos_x': e.pos_x, 'pos_y': e.pos_y, 'width': e.width, 'height': e.height, 'content': e.content, 'font_size': e.font_size} 
            for e in [title_element, subtitle_element]
        ]
    }
    
    return jsonify({
        'id': new_presentation.id,
        'title': new_presentation.title,
        'updated_at': new_presentation.updated_at.isoformat(),
        'first_slide': first_slide_data
    }), 201

@presentations_bp.route('/presentations/<string:presentation_id>', methods=['GET'])
@token_required
def get_presentation_by_id(presentation_id):
    presentation = Presentation.query.get_or_404(presentation_id)
    if presentation.user_id != g.current_user.id: return jsonify({'message': 'Доступ запрещен'}), 403
    
    slides_output = []
    slides = Slide.query.filter_by(presentation_id=presentation.id).order_by(Slide.slide_number).all()
    for slide in slides:
        elements = SlideElement.query.filter_by(slide_id=slide.id).all()
        elements_output = [{'id': e.id, 'element_type': e.element_type, 'pos_x': e.pos_x, 'pos_y': e.pos_y, 'width': e.width, 'height': e.height, 'content': e.content, 'font_size': e.font_size} for e in elements]
        slides_output.append({'id': slide.id, 'slide_number': slide.slide_number, 'background_color': slide.background_color, 'elements': elements_output})
        
    return jsonify({'id': presentation.id, 'title': presentation.title, 'slides': slides_output}), 200

@presentations_bp.route('/presentations', methods=['GET'])
@token_required
def get_presentations():
    presentations = Presentation.query.filter_by(user_id=g.current_user.id).order_by(Presentation.updated_at.desc()).all()
    output = []
    for p in presentations:
        first_slide = Slide.query.filter_by(presentation_id=p.id, slide_number=1).first()
        if first_slide:
            elements = SlideElement.query.filter_by(slide_id=first_slide.id).all()
            first_slide_data = {
                'id': first_slide.id,
                'slide_number': first_slide.slide_number,
                'background_color': first_slide.background_color,
                'elements': [{'id': e.id, 'element_type': e.element_type, 'pos_x': e.pos_x, 'pos_y': e.pos_y, 'width': e.width, 'height': e.height, 'content': e.content, 'font_size': e.font_size} for e in elements]
            }
        else:
            first_slide_data = None
        
        output.append({
            'id': p.id,
            'title': p.title,
            'updated_at': p.updated_at.isoformat(),
            'first_slide': first_slide_data
        })
    return jsonify(output), 200

@presentations_bp.route('/presentations/<string:presentation_id>', methods=['DELETE'])
@token_required
def delete_presentation(presentation_id):
    presentation = Presentation.query.get_or_404(presentation_id)
    if presentation.user_id != g.current_user.id: return jsonify({'message': 'Доступ запрещен'}), 403
    db.session.delete(presentation)
    db.session.commit()
    return jsonify({'message': 'Презентация успешно удалена'}), 200

@presentations_bp.route('/presentations/<string:presentation_id>', methods=['PUT'])
@token_required
def update_presentation(presentation_id):
    presentation = Presentation.query.get_or_404(presentation_id)
    if presentation.user_id != g.current_user.id: return jsonify({'message': 'Доступ запрещен'}), 403
    data = request.get_json()
    if 'title' in data: presentation.title = data['title']
    db.session.commit()
    
    first_slide = Slide.query.filter_by(presentation_id=presentation.id, slide_number=1).first()
    if first_slide:
        elements = SlideElement.query.filter_by(slide_id=first_slide.id).all()
        first_slide_data = {
            'id': first_slide.id,
            'slide_number': first_slide.slide_number,
            'background_color': first_slide.background_color,
            'elements': [{'id': e.id, 'element_type': e.element_type, 'pos_x': e.pos_x, 'pos_y': e.pos_y, 'width': e.width, 'height': e.height, 'content': e.content, 'font_size': e.font_size} for e in elements]
        }
    else:
        first_slide_data = None

    return jsonify({
        'id': presentation.id,
        'title': presentation.title,
        'updated_at': presentation.updated_at.isoformat(),
        'first_slide': first_slide_data
    }), 200